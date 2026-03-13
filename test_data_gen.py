import os
from pptx import Presentation
from openpyxl import Workbook

def generate_test_files():
    os.makedirs('test_samples', exist_ok=True)
    
    # Generate PPTX
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = width = height = 1000000
    
    # 2 copies of "This is a very specific confidential string for testing."
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is a very specific confidential string for testing."
    
    txBox2 = slide.shapes.add_textbox(left, top + 1000000, width, height)
    tf2 = txBox2.text_frame
    tf2.text = "This is a very specific confidential string for testing."

    prs.save('test_samples/sample.pptx')
    
    # Generate XLSX
    wb = Workbook()
    ws = wb.active
    
    # 3 copies of "Another highly secret project delta phrase."
    ws['A1'] = "Another highly secret project delta phrase."
    ws['B2'] = "Another highly secret project delta phrase."
    ws['C3'] = "Another highly secret project delta phrase."
    
    # And 1 copy of the string from PPTX to verify cross-file counting
    ws['D4'] = "This is a very specific confidential string for testing."

    wb.save('test_samples/sample.xlsx')

    print("Test files generated successfully in test_samples/")

if __name__ == "__main__":
    generate_test_files()
