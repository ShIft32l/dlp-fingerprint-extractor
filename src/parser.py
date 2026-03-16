import logging
from pathlib import Path
from typing import Optional, Union, List
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation

try:
    import pytesseract
    from pdf2image import convert_from_path
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

logger = logging.getLogger(__name__)

def read_txt(path: Path) -> Optional[str]:
    """Reads raw text from a .txt file."""
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        logger.error(f"Error reading TXT file {path}: {e}")
        return None

def read_docx(path: Path) -> Optional[str]:
    """Reads text from a MS Word .docx file."""
    try:
        doc = Document(str(path))
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return "\n".join(full_text)
    except Exception as e:
        logger.error(f"Error reading DOCX file {path}: {e}")
        return None

def _ocr_pdf_pages(path: Path, page_indices: List[int]) -> List[str]:
    """OCR specific pages of a PDF using pytesseract.
    
    Args:
        path: Path to the PDF file.
        page_indices: 0-based indices of pages to OCR.
    
    Returns:
        List of extracted text strings, one per page.
    """
    ocr_texts = []
    try:
        # Convert only the needed pages to images
        # pdf2image uses 1-based page numbers
        for idx in page_indices:
            images = convert_from_path(
                str(path),
                first_page=idx + 1,
                last_page=idx + 1,
                dpi=300
            )
            if images:
                text = pytesseract.image_to_string(images[0])
                ocr_texts.append(text.strip() if text else "")
            else:
                ocr_texts.append("")
    except Exception as e:
        logger.error(f"OCR failed for {path}: {e}")
        ocr_texts.extend([""] * (len(page_indices) - len(ocr_texts)))
    return ocr_texts

def read_pdf(path: Path) -> Optional[str]:
    """Reads text from a PDF file.
    
    Uses a two-pass strategy:
    1. Try PyPDF2 text extraction (fast).
    2. For pages with no text, fall back to OCR via pytesseract.
    """
    try:
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            num_pages = len(reader.pages)
            text_blocks = [""] * num_pages
            ocr_needed = []

            # Pass 1: PyPDF2 text extraction
            for i, page in enumerate(reader.pages):
                extracted = page.extract_text()
                if extracted and extracted.strip():
                    text_blocks[i] = extracted
                else:
                    ocr_needed.append(i)

        # Pass 2: OCR fallback for empty pages
        if ocr_needed:
            if not OCR_AVAILABLE:
                logger.warning(
                    f"PDF '{path.name}' has {len(ocr_needed)} image-based page(s) "
                    f"but OCR libs (pytesseract, pdf2image) are not installed. "
                    f"Install them with: pip install pytesseract pdf2image Pillow"
                )
            else:
                logger.info(
                    f"Running OCR on {len(ocr_needed)}/{num_pages} page(s) of '{path.name}'..."
                )
                ocr_results = _ocr_pdf_pages(path, ocr_needed)
                for idx, ocr_text in zip(ocr_needed, ocr_results):
                    if ocr_text:
                        text_blocks[idx] = ocr_text

        # Filter out empty blocks and join
        result = "\n".join(block for block in text_blocks if block.strip())
        return result if result.strip() else None

    except Exception as e:
        logger.error(f"Error reading PDF file {path}: {e}")
        return None

def read_xlsx(path: Path) -> Optional[str]:
    """Reads text from an Excel .xlsx file."""
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        text_blocks = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None and str(cell).strip():
                        text_blocks.append(str(cell).strip())
        return "\n".join(text_blocks)
    except Exception as e:
        logger.error(f"Error reading XLSX file {path}: {e}")
        return None

def read_pptx(path: Path) -> Optional[str]:
    """Reads text from a PowerPoint .pptx file."""
    try:
        prs = Presentation(str(path))
        text_blocks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_blocks.append(shape.text.strip())
        return "\n".join(text_blocks)
    except Exception as e:
        logger.error(f"Error reading PPTX file {path}: {e}")
        return None

def extract_document_text(path: Union[str, Path]) -> Optional[str]:
    """Factory function to extract text based on file suffix."""
    file_path = Path(path)
    
    if not file_path.exists():
        logger.warning(f"File not found: {file_path}")
        return None
    
    if not file_path.is_file():
        logger.warning(f"Path is not a regular file: {file_path}")
        return None

    suffix = file_path.suffix.lower()
    
    if suffix == ".txt":
        return read_txt(file_path)
    elif suffix == ".docx":
        return read_docx(file_path)
    elif suffix == ".pdf":
        return read_pdf(file_path)
    elif suffix == ".xlsx":
        return read_xlsx(file_path)
    elif suffix == ".pptx":
        return read_pptx(file_path)
    else:
        logger.warning(f"Unsupported file format '{suffix}' for file: {file_path}")
        return None
